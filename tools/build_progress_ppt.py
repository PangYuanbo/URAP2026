from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
from copy import deepcopy

TEMPLATE = r"C:\Users\aaron\Desktop\URAP\URAP-UAV-to-UAV-Detection-and-Tracking\doc\Copy of Moving Object Detection from Moving Platform.pptx"
OUTPUT = r"C:\Users\aaron\Desktop\URAP\URAP-UAV-to-UAV-Detection-and-Tracking\doc\URAP_Progress_Update_2026-03-23.pptx"

PURPLE = RGBColor(164, 146, 232)
PURPLE_DARK = RGBColor(110, 92, 194)
PURPLE_LIGHT = RGBColor(243, 239, 255)
TEXT = RGBColor(35, 31, 32)
GREEN = RGBColor(64, 136, 102)
ORANGE = RGBColor(212, 136, 42)
RED = RGBColor(187, 67, 67)
GRAY = RGBColor(96, 96, 96)

prs = Presentation(TEMPLATE)

# Remove all template slides while keeping master/layout/theme.
slides = prs.slides._sldIdLst
for i in range(len(slides) - 1, -1, -1):
    rId = slides[i].rId
    prs.part.drop_rel(rId)
    del slides[i]

layouts = {layout.name: layout for layout in prs.slide_layouts}


def set_run_style(run, size=20, bold=False, color=TEXT):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = 'Aptos'


def add_title_slide(title, subtitle):
    slide = prs.slides.add_slide(layouts['TITLE'])
    slide.shapes.title.text = title
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle
    return slide


def add_section_header(title, subtitle=''):
    slide = prs.slides.add_slide(layouts['SECTION_HEADER'])
    slide.shapes.title.text = title
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle
    return slide


def clear_placeholder(ph):
    ph.text = ''
    tf = ph.text_frame
    tf.clear()
    return tf


def add_bullets_slide(title, bullets, footer=None):
    slide = prs.slides.add_slide(layouts['TITLE_AND_BODY'])
    slide.shapes.title.text = title
    body = slide.placeholders[1]
    tf = clear_placeholder(body)
    tf.word_wrap = True
    first = True
    for item in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        p.level = level
        p.text = text
        for r in p.runs:
            set_run_style(r, size=22 if level == 0 else 18)
    if footer:
        tx = slide.shapes.add_textbox(Inches(0.6), Inches(6.85), Inches(12.0), Inches(0.35))
        p = tx.text_frame.paragraphs[0]
        p.text = footer
        for r in p.runs:
            set_run_style(r, size=11, color=GRAY)
    return slide


def add_two_column_slide(title, left_title, left_bullets, right_title, right_bullets, footer=None):
    slide = prs.slides.add_slide(layouts['TITLE_ONLY'])
    slide.shapes.title.text = title

    def add_panel(x, y, w, h, head, bullets):
        rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        rect.fill.solid()
        rect.fill.fore_color.rgb = PURPLE_LIGHT
        rect.line.color.rgb = PURPLE
        rect.line.width = Pt(1.6)
        tx = rect.text_frame
        tx.clear()
        tx.word_wrap = True
        p = tx.paragraphs[0]
        p.text = head
        for r in p.runs:
            set_run_style(r, size=22, bold=True, color=PURPLE_DARK)
        for item in bullets:
            p = tx.add_paragraph()
            if isinstance(item, tuple):
                text, level = item
            else:
                text, level = item, 0
            p.level = level
            p.text = text
            for r in p.runs:
                set_run_style(r, size=18 if level == 0 else 16)
        return rect

    add_panel(Inches(0.55), Inches(1.45), Inches(5.85), Inches(4.95), left_title, left_bullets)
    add_panel(Inches(6.7), Inches(1.45), Inches(5.85), Inches(4.95), right_title, right_bullets)

    if footer:
        tx = slide.shapes.add_textbox(Inches(0.6), Inches(6.65), Inches(12.0), Inches(0.4))
        p = tx.text_frame.paragraphs[0]
        p.text = footer
        for r in p.runs:
            set_run_style(r, size=11, color=GRAY)
    return slide


def add_kpi_slide():
    slide = prs.slides.add_slide(layouts['BLANK'])
    title = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(10.5), Inches(0.6))
    p = title.text_frame.paragraphs[0]
    p.text = 'Current Reproduction Status and Core Findings'
    for r in p.runs:
        set_run_style(r, size=28, bold=True)

    cards = [
        ('TransVisDrone', ['NPS val/test reproduced', 'NPS val mAP@0.5 = 0.948', 'NPS test mAP@0.5 = 0.938', 'AOT fulltest recomputation finished'], GREEN),
        ('AOT Winner v022', ['AOT fulltest 172/172 flights completed', 'HFAR ~ 0.523, EDR@300 ~ 0.989', 'Very strong on AOT, weak on NPS generalization'], ORANGE),
        ('ESOD / YOLOMG Insights', ['ESOD reproduced on VisDrone', 'YOLOMG useful as motion-prior donor', 'Do not use full YOLOMG as backbone'], PURPLE_DARK),
    ]
    xs = [0.65, 4.35, 8.05]
    for (title_txt, bullet_list, color), x in zip(cards, xs):
        rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.5), Inches(3.2), Inches(3.6))
        rect.fill.solid(); rect.fill.fore_color.rgb = RGBColor(250,250,252)
        rect.line.color.rgb = color; rect.line.width = Pt(2.2)
        tf = rect.text_frame; tf.clear(); tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = title_txt
        for r in p.runs: set_run_style(r, size=22, bold=True, color=color)
        for b in bullet_list:
            p = tf.add_paragraph(); p.text = b; p.level = 0
            for r in p.runs: set_run_style(r, size=16)

    note = slide.shapes.add_textbox(Inches(0.75), Inches(5.45), Inches(11.6), Inches(1.2))
    p = note.text_frame.paragraphs[0]
    p.text = 'Main engineering decision: keep TransVisDrone as the only detector backbone; transfer only the reusable low-FP and motion-prior ideas from Winner / YOLOMG / ESOD.'
    for r in p.runs:
        set_run_style(r, size=18, color=PURPLE_DARK)
    return slide


def add_architecture_slide():
    slide = prs.slides.add_slide(layouts['BLANK'])
    title = slide.shapes.add_textbox(Inches(0.6), Inches(0.25), Inches(11.8), Inches(0.6))
    p = title.text_frame.paragraphs[0]
    p.text = 'Planned V1 Architecture: TVD Backbone + Motion Prior + Selective Refinement'
    for r in p.runs:
        set_run_style(r, size=26, bold=True)

    y = Inches(2.2)
    boxes = [
        ('Input Video', 0.5, 1.3),
        ('CMC', 1.9, 1.0),
        ('Motion Prior\n(YOLOMG donor only)', 3.1, 1.8),
        ('TransVisDrone\nMain detector', 5.3, 1.7),
        ('Score Fusion /\nRe-score', 7.4, 1.3),
        ('HBS', 9.0, 0.9),
        ('Selective Refinement\nOEN + scheduler + ESOD + EDTC', 10.2, 2.0),
    ]
    created = []
    for text, x, w in boxes:
        rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), y, Inches(w), Inches(1.0))
        rect.fill.solid(); rect.fill.fore_color.rgb = PURPLE_LIGHT
        rect.line.color.rgb = PURPLE; rect.line.width = Pt(1.8)
        tf = rect.text_frame; tf.clear(); tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; p.text = text
        for r in p.runs: set_run_style(r, size=16, bold=True)
        created.append(rect)

    for a, b in zip(created[:-1], created[1:]):
        x1 = a.left + a.width
        y1 = a.top + a.height/2
        x2 = b.left
        y2 = b.top + b.height/2
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
        line.line.color.rgb = PURPLE_DARK
        line.line.width = Pt(2)

    confirm = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.7), Inches(4.25), Inches(2.0), Inches(0.9))
    confirm.fill.solid(); confirm.fill.fore_color.rgb = RGBColor(246,248,250)
    confirm.line.color.rgb = PURPLE_DARK; confirm.line.width = Pt(1.8)
    p = confirm.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; p.text = 'confirm-lite'
    for r in p.runs: set_run_style(r, size=17, bold=True)
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, created[-1].left + created[-1].width/2, created[-1].top + created[-1].height, confirm.left + confirm.width/2, confirm.top)
    line.line.color.rgb = PURPLE_DARK; line.line.width = Pt(1.8)

    final = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.2), Inches(4.25), Inches(2.2), Inches(0.9))
    final.fill.solid(); final.fill.fore_color.rgb = RGBColor(251,251,251)
    final.line.color.rgb = GREEN; final.line.width = Pt(1.8)
    p = final.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; p.text = 'Final outputs'
    for r in p.runs: set_run_style(r, size=17, bold=True, color=GREEN)
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, confirm.left + confirm.width, confirm.top + confirm.height/2, final.left, final.top + final.height/2)
    line.line.color.rgb = PURPLE_DARK; line.line.width = Pt(1.8)

    note = slide.shapes.add_textbox(Inches(0.8), Inches(5.55), Inches(11.4), Inches(1.1))
    p = note.text_frame.paragraphs[0]
    p.text = 'Key rule: YOLOMG is not used as a second detector. Only its motion-cue idea is kept, then fused with TVD detections.'
    for r in p.runs: set_run_style(r, size=17, color=PURPLE_DARK)
    return slide


def add_module_roles_slide():
    slide = prs.slides.add_slide(layouts['BLANK'])
    title = slide.shapes.add_textbox(Inches(0.6), Inches(0.25), Inches(11.8), Inches(0.6))
    p = title.text_frame.paragraphs[0]; p.text = 'Module Roles in the Current Plan'
    for r in p.runs: set_run_style(r, size=28, bold=True)

    rows = [
        ('TransVisDrone', 'Main video detector backbone', 'Primary detector'),
        ('CMC', 'Remove ego-motion before trusting temporal cues', 'Pre-processing'),
        ('YOLOMG donor', 'Generate motion prior / difference-map evidence', 'Auxiliary evidence only'),
        ('HBS', 'Suppress urban structured clutter (buildings / edges / wires)', 'Precision / low-FP block'),
        ('GPT-5', 'Research synthesis, literature screening, ablation planning', 'Offline research tool'),
        ('Typeless', 'Fast implementation / iteration support layer', 'Engineering acceleration, not runtime inference'),
    ]
    table = slide.shapes.add_table(len(rows)+1, 3, Inches(0.55), Inches(1.35), Inches(12.0), Inches(4.5)).table
    headers = ['Component', 'What It Does', 'Status in the System']
    for c, h in enumerate(headers):
        cell = table.cell(0,c); cell.text = h
        for p in cell.text_frame.paragraphs:
            for r in p.runs: set_run_style(r, size=16, bold=True, color=PURPLE_DARK)
        cell.fill.solid(); cell.fill.fore_color.rgb = PURPLE_LIGHT
    for r_idx, row in enumerate(rows, start=1):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx, c_idx); cell.text = val
            for p in cell.text_frame.paragraphs:
                for r in p.runs: set_run_style(r, size=14)
    note = slide.shapes.add_textbox(Inches(0.8), Inches(6.05), Inches(11.2), Inches(0.6))
    p = note.text_frame.paragraphs[0]
    p.text = 'Only TransVisDrone, CMC, motion prior, HBS, refinement and confirm-lite belong to the runtime model. GPT-5 / Typeless stay outside the onboard inference path.'
    for r in p.runs: set_run_style(r, size=16, color=GRAY)
    return slide


add_title_slide(
    'URAP Progress Update',
    'Toward edge-deployable tiny-UAV / tiny-obstacle detection in urban clutter | 2026-03-23'
)
add_kpi_slide()
add_bullets_slide(
    'TransVisDrone: What the Backbone Actually Is',
    [
        'Spatial branch: CSPDarkNet-53 style CNN extracts per-frame spatial features.',
        'Temporal branch: Video Swin learns short-range spatio-temporal dependencies from video clips.',
        'Fusion: spatial + temporal evidence are fused before final detection.',
        'Strengths: native video detector, good NPS reproduction, extensible to AOT.',
        'Limitations: still too weak on AOT false-alarm control without extra engineering.',
        'Why we keep it: it is a clean, neutral platform backbone rather than a highly specialized heuristic stack.',
    ],
    footer='Source base: C:\\Users\\aaron\\Desktop\\URAP\\URAP-UAV-to-UAV-Detection-and-Tracking\\doc\\transvisdrone_method_explained.md'
)
add_bullets_slide(
    'YOLOMG / Image-First Motion Line: What We Keep and What We Do Not Keep',
    [
        'Core idea: use compensated frame-difference / mask-like motion cues to make tiny moving objects more visible.',
        'Strength: motion cue is highly useful when the object is too small to classify confidently from appearance alone.',
        'Weakness: the full YOLOMG pipeline depends strongly on its own second stream and data construction.',
        'Observed in our work: cross-domain generalization is much weaker than TVD when moved out of its native setting.',
        'Decision: do not use full YOLOMG as backbone.',
        'Keep only the donor idea: motion-prior generation after CMC, then fuse with TVD detections.',
    ],
    footer='Related evidence: C:\\Users\\aaron\\Desktop\\URAP\\URAP-UAV-to-UAV-Detection-and-Tracking\\doc\\repro_yolomg_ard100.md and generalization notes'
)
add_two_column_slide(
    'Backbone Decision: TransVisDrone vs Image-Only / YOLOMG-style Line',
    'Why TransVisDrone stays the backbone',
    [
        'Native video detector with explicit temporal modeling.',
        'Cleaner integration point for low-FP post-processing, clutter suppression and refinement.',
        'Already reproduced on NPS and fully evaluated on AOT.',
        'More neutral platform for adding new modules without rewriting the entire detector logic.',
    ],
    'Why YOLOMG is not the backbone',
    [
        'Its motion branch is deeply tied to its own detection logic.',
        'If used as backbone, motion reasoning gets duplicated with our later temporal modules.',
        'Observed generalization weakness makes it risky as the central skeleton.',
        'Best role is motion-prior donor, not full detector replacement.',
    ],
    footer='Assumption used here: “Unimage” in the request is treated as the image-first / single-image line currently represented by YOLOMG-style design choices.'
)
add_architecture_slide()
add_module_roles_slide()
add_bullets_slide(
    'How the Framework Should Move Forward',
    [
        'Phase 1: freeze the detector skeleton -> TVD + CMC-guided motion prior + confirm-lite.',
        'Phase 2: add only one clutter suppressor winner -> HBS-lite.',
        'Phase 3: add one bounded selective-refinement block -> OEN + scheduler + ESOD + EDTC.',
        'Phase 4: run ablations block-by-block on both NPS and AOT; do not add same-slot comparators to the mainline.',
        'Phase 5: only after V1 is stable, consider V2 items such as stronger ROI logic, structural small-object surgery, or RL trigger policy.',
    ],
    footer='Design rule: one slot, one winner. Extra modules move to V2 / appendix / ablation-only.'
)
add_bullets_slide(
    'Immediate Next Steps',
    [
        '1. Freeze V1 architecture and stop expanding the mainline search space.',
        '2. Implement the CMC + motion-prior donor path around the existing TVD inference code.',
        '3. Integrate HBS as the only clutter-suppression block in V1.',
        '4. Keep GPT-5 / Typeless on the research and implementation side, not inside runtime inference.',
        '5. Re-run AOT and NPS with strict ablation bookkeeping and decide whether the selective refinement block is justified in V1 or should be deferred to V1.1.',
    ],
    footer='Prepared from current project docs, reproduction logs, and architecture review notes in the local workspace.'
)

prs.save(OUTPUT)
print(OUTPUT)
print('slides', len(list(prs.slides)))
