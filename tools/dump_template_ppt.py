from pptx import Presentation
ppt = r"C:\Users\aaron\Desktop\URAP\URAP-UAV-to-UAV-Detection-and-Tracking\doc\Copy of Moving Object Detection from Moving Platform.pptx"
p = Presentation(ppt)
for i, s in enumerate(list(p.slides), 1):
    print(f'===== SLIDE {i} | layout={s.slide_layout.name} =====')
    for j, sh in enumerate(s.shapes):
        txt = getattr(sh, 'text', '').strip().replace('\n', ' | ')
        if txt:
            print(f'[{j}] {txt[:800]}')
