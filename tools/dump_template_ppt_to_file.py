from pptx import Presentation
ppt = r"C:\Users\aaron\Desktop\URAP\URAP-UAV-to-UAV-Detection-and-Tracking\doc\Copy of Moving Object Detection from Moving Platform.pptx"
out = r"C:\Users\aaron\Desktop\URAP\URAP-UAV-to-UAV-Detection-and-Tracking\doc\template_ppt_dump_utf8.txt"
p = Presentation(ppt)
with open(out, 'w', encoding='utf-8') as f:
    for i, s in enumerate(list(p.slides), 1):
        f.write(f'===== SLIDE {i} | layout={s.slide_layout.name} =====\n')
        for j, sh in enumerate(s.shapes):
            txt = getattr(sh, 'text', '').strip().replace('\n', ' | ')
            if txt:
                f.write(f'[{j}] {txt[:2000]}\n')
print(out)
