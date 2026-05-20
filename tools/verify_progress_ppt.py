from pptx import Presentation
p = Presentation(r"C:\Users\aaron\Desktop\URAP\URAP-UAV-to-UAV-Detection-and-Tracking\doc\URAP_Progress_Update_2026-03-23.pptx")
print('slides', len(list(p.slides)))
for i,s in enumerate(list(p.slides),1):
    title = ''
    for sh in s.shapes:
        if getattr(sh, 'has_text_frame', False):
            txt = sh.text.strip().replace('\n',' | ')
            if txt:
                title = txt[:140]
                break
    print(i, title)
